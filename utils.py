import hashlib
from aiogram import BaseMiddleware
from aiogram.types import Message
from typing import Callable, Dict, Any, Awaitable
from cachetools import TTLCache
import os
import time
import re
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, Cm
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt

# We use a simple TTL cache to keep track of requests per user ID
# 1-second TTL means 1 message per second allowed per user
rate_limit_cache = TTLCache(maxsize=10000, ttl=1)

def hash_password(password: str) -> str:
    """Returns SHA-256 hash of a given password."""
    return hashlib.sha256(password.encode()).hexdigest()

async def send_split_message(message: Message, text: str):
    """Splits a long message into chunks and sends them sequentially."""
    limit = 4000  # Telegram limit is 4096, using 4000 for safety
    if len(text) <= limit:
        await message.answer(text)
    else:
        for i in range(0, len(text), limit):
            chunk = text[i:i + limit]
            await message.answer(chunk)

def clean_markdown(text: str) -> str:
    """Removes common markdown-style formatting symbols."""
    # Remove Bold/Italic stars
    text = re.sub(r'\*+', '', text)
    # Remove Hashtags headings (###, ####, etc.)
    text = re.sub(r'#+\s*', '', text)
    # Remove Underscores
    text = re.sub(r'_+', '', text)
    # Remove horizontal lines represented by dashes
    text = re.sub(r'-{3,}', '', text)
    return text.strip()

def cleanup_old_files():
    """Deletes files older than 24 hours from the exports directory."""
    exports_dir = "exports"
    if not os.path.exists(exports_dir):
        return
    
    now = time.time()
    for f in os.listdir(exports_dir):
        f_path = os.path.join(exports_dir, f)
        if os.path.isfile(f_path):
            # 86400 seconds = 24 hours
            if os.stat(f_path).st_mtime < now - 86400:
                try:
                    os.remove(f_path)
                except Exception:
                    pass

def create_academic_docx(text: str, filename: str, topic: str, subject: str, doc_type: str, author: str) -> str:
    """Generates a structured .docx file with Cover, Outline, and Content."""
    exports_dir = "exports"
    if not os.path.exists(exports_dir):
        os.makedirs(exports_dir)
    
    cleanup_old_files()
    
    doc = Document()
    
    # 1. Page Setup (Margins 2cm)
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2)
        section.right_margin = Cm(2)
        
    # 2. Configure Styles
    # Normal (Body Text)
    style_normal = doc.styles['Normal']
    font_normal = style_normal.font
    font_normal.name = 'Times New Roman'
    font_normal.size = Pt(14)
    font_normal.color.rgb = None # Automatic/Black
    style_normal.paragraph_format.line_spacing = 1.5
    style_normal.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    style_normal.paragraph_format.space_after = Pt(6) # Subtle spacing between paragraphs
    
    # Heading 1 (Main sections)
    style_h1 = doc.styles['Heading 1']
    font_h1 = style_h1.font
    font_h1.name = 'Times New Roman'
    font_h1.size = Pt(16)
    font_h1.bold = True
    font_h1.color.rgb = None
    style_h1.paragraph_format.line_spacing = 1.5
    style_h1.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER # Headings centered
    style_h1.paragraph_format.space_before = Pt(12)
    style_h1.paragraph_format.space_after = Pt(6)
    style_h1.paragraph_format.keep_with_next = True
    
    # Heading 2 (Subsections)
    style_h2 = doc.styles['Heading 2']
    font_h2 = style_h2.font
    font_h2.name = 'Times New Roman'
    font_h2.size = Pt(14)
    font_h2.bold = True
    font_h2.color.rgb = None
    style_h2.paragraph_format.line_spacing = 1.5
    style_h2.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
    style_h2.paragraph_format.space_before = Pt(12)
    style_h2.paragraph_format.space_after = Pt(6)
    style_h2.paragraph_format.keep_with_next = True

    # --- COVER PAGE ---
    # Institution Name
    p_gov = doc.add_paragraph()
    p_gov.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_gov = p_gov.add_run(
        "O'ZBEKISTON RESPUBLIKASI OLIY TA'LIM, FAN VA INNOVATSIYALAR VAZIRLIGI\n"
        "__________________________________________________________________"
    )
    run_gov.bold = True
    run_gov.font.size = Pt(12)
    
    for _ in range(4): doc.add_paragraph("")
    
    # Subject/Fan
    p_fac = doc.add_paragraph()
    p_fac.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r_fac = p_fac.add_run(f"FAN: {subject.upper()}")
    r_fac.bold = True
    r_fac.font.size = Pt(14)
    
    # Document Type
    p_type = doc.add_paragraph()
    p_type.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r_type = p_type.add_run(doc_type.upper())
    r_type.bold = True
    r_type.font.size = Pt(22)
    
    # Topic
    p_top = doc.add_paragraph()
    p_top.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r_top = p_top.add_run(f"MAVZU: {topic.upper()}")
    r_top.bold = True
    r_top.font.size = Pt(14)
    
    for _ in range(4): doc.add_paragraph("")
    
    # Signatures
    p_sign = doc.add_paragraph()
    p_sign.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    r_sign = p_sign.add_run(
        f"Bajardi: {author}\n"
        f"Tekshirdi: _____________________"
    )
    r_sign.font.size = Pt(14)
    
    for _ in range(4): doc.add_paragraph("")
    
    # Bottom Year
    p_year = doc.add_paragraph()
    p_year.alignment = WD_ALIGN_PARAGRAPH.CENTER
    import datetime
    year = datetime.datetime.now().year
    p_year.add_run(f"Toshkent - {year}")
    
    doc.add_page_break()
    
    # --- MUNDARIJA (TABLE OF CONTENTS) ---
    p_toc = doc.add_paragraph()
    p_toc.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_toc.add_run("MUNDARIJA").bold = True
    
    # Embed Table of Contents
    paragraph = doc.add_paragraph()
    run = paragraph.add_run()
    fldChar = OxmlElement('w:fldChar')
    fldChar.set(qn('w:fldCharType'), 'begin')
    instrText = OxmlElement('w:instrText')
    instrText.set(qn('xml:space'), 'preserve')
    instrText.text = 'TOC \\o "1-3" \\h \\z \\u'
    fldChar2 = OxmlElement('w:fldChar')
    fldChar2.set(qn('w:fldCharType'), 'separate')
    fldChar3 = OxmlElement('w:fldChar')
    fldChar3.set(qn('w:fldCharType'), 'end')
    r = run._r
    r.append(fldChar)
    r.append(instrText)
    r.append(fldChar2)
    r.append(fldChar3)
    
    doc.add_page_break()
    
    # --- CONTENT PARSING & GENERATION ---
    content_parts = text.split("REJA:", 1)
    if len(content_parts) > 1:
        parts = content_parts[1].split("\n\n", 1)
        actual_content = parts[1].strip() if len(parts) > 1 else parts[0].strip()
    else:
        re_match = re.split(r'KIRISH|Kirish', text, maxsplit=1)
        if len(re_match) > 1:
            actual_content = "KIRISH\n" + re_match[1].strip()
        else:
            actual_content = text
            
    paragraphs = actual_content.split('\n')
    
    first_heading = True
    for line in paragraphs:
        if not line.strip():
            continue
            
        line_clean = clean_markdown(line)
        
        # Check if the line is a main heading (Heading 1)
        is_h1 = False
        h1_keywords = ["KIRISH", "XULOSA", "ASOSIY QISM", "FOYDALANILGAN ADABIYOTLAR", "I BOB:", "II BOB:", "III BOB:", "REJA:"]
        for kw in h1_keywords:
            if line_clean.upper().startswith(kw):
                is_h1 = True
                break
                
        # Check if the line is a sub-heading (Heading 2)
        is_h2 = False
        if not is_h1:
            if re.match(r'^\d+(\.\d+)+\.?\s+', line_clean) or re.match(r'^\d+\.\s+', line_clean):
                is_h2 = True
                
        if is_h1:
            if not first_heading:
                doc.add_page_break()
            else:
                first_heading = False
                
            p = doc.add_paragraph(line_clean, style='Heading 1')
        elif is_h2:
            p = doc.add_paragraph(line_clean, style='Heading 2')
        else:
            p = doc.add_paragraph()
            if line_clean.isupper() and len(line_clean) > 5 and len(line_clean) < 100:
                run = p.add_run(line_clean)
                run.bold = True
            else:
                p.add_run(line_clean)
                
    file_path = os.path.abspath(os.path.join(exports_dir, filename))
    doc.save(file_path)
    return file_path

# Keep wrapper functions for compatibility with other files if needed
def create_docx(text: str, filename: str, university: str, author: str, topic: str, doc_type: str) -> str:
    return create_academic_docx(text, filename, topic, university, doc_type, author)

def create_mustaqil_ish_docx(text: str, filename: str, topic: str, subject: str, university: str, teacher: str, author: str) -> str:
    return create_academic_docx(text, filename, topic, subject, "MUSTAQIL ISH", author)

def create_pptx(text: str, filename: str, topic: str, author: str) -> str:
    """Generates a styled, modern .pptx file with slides from structured text."""
    exports_dir = "exports"
    if not os.path.exists(exports_dir):
        os.makedirs(exports_dir)
    
    cleanup_old_files()
    
    prs = Presentation()
    
    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    from pptx.dml.color import RGBColor
    # RGB Colors
    BG_COLOR = RGBColor(248, 250, 252) # Slate 50
    TITLE_COLOR = RGBColor(15, 23, 42) # Slate 900
    BODY_COLOR = RGBColor(71, 85, 105) # Slate 600
    
    # --- SLIDE 1: TITLE SLIDE ---
    title_slide_layout = prs.slide_layouts[5] # Blank
    slide = prs.slides.add_slide(title_slide_layout)
    
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(15, 23, 42) # Dark background for title
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = topic.upper()
    p.font.name = 'Arial'
    p.font.size = PptxPt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = f"\nTaqdimotchi: {author}\n2025-2026 O'quv Yili"
    p2.font.name = 'Arial'
    p2.font.size = PptxPt(20)
    p2.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # --- SLIDES 2+: CONTENT ---
    sections = re.split(r'Slayd \d+:', text)
    if len(sections) == 1:
        sections = text.split('\n\n')
        
    for section in sections:
        section = section.strip()
        if not section:
            continue
            
        lines = [l.strip() for l in section.split('\n') if l.strip()]
        if not lines:
            continue
            
        slide_title = lines[0]
        slide_title = clean_markdown(slide_title)
        slide_body_lines = lines[1:]
        
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank slide
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG_COLOR
        
        titleBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.333), Inches(1))
        tf_title = titleBox.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_title
        p_title.font.name = 'Arial'
        p_title.font.size = PptxPt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = TITLE_COLOR
        
        bodyBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(4.5))
        tf_body = bodyBox.text_frame
        tf_body.word_wrap = True
        
        first_bullet = True
        for line in slide_body_lines:
            line_clean = clean_markdown(line)
            if line_clean.startswith('-'):
                line_clean = line_clean.lstrip('-').strip()
            elif line_clean.startswith('*'):
                line_clean = line_clean.lstrip('*').strip()
                
            if not line_clean:
                continue
                
            if first_bullet:
                p_bullet = tf_body.paragraphs[0]
                first_bullet = False
            else:
                p_bullet = tf_body.add_paragraph()
                
            p_bullet.text = f"•  {line_clean}"
            p_bullet.font.name = 'Arial'
            p_bullet.font.size = PptxPt(18)
            p_bullet.font.color.rgb = BODY_COLOR
            p_bullet.space_after = PptxPt(12)

    file_path = os.path.abspath(os.path.join(exports_dir, filename))
    prs.save(file_path)
    return file_path

class ThrottlingMiddleware(BaseMiddleware):
    """Simple anti-spam middleware to prevent rapid messaging (Rate Limit)."""
    async def __call__(
        self,
        handler: Callable[[Message, Dict[str, Any]], Awaitable[Any]],
        event: Message,
        data: Dict[str, Any]
    ) -> Any:
        if event.from_user:
            user_id = event.from_user.id
            if user_id in rate_limit_cache:
                # Silently drop the spamming message
                return
            rate_limit_cache[user_id] = True
        return await handler(event, data)

class BannedUserMiddleware(BaseMiddleware):
    """Middleware to check if the user is banned from using the bot."""
    async def __call__(
        self,
        handler: Callable[[Message, Dict[str, Any]], Awaitable[Any]],
        event: Message,
        data: Dict[str, Any]
    ) -> Any:
        from database import is_banned
        if event.from_user:
            if await is_banned(event.from_user.id):
                # Optionally inform the user they are banned, or just drop
                return
        return await handler(event, data)

def parse_pages_input(service_type: str, text: str) -> str:
    """Parses user text input to map it to a valid pricing tier."""
    import re
    match = re.search(r'(\d+)', text)
    if not match:
        if service_type == "taqdimot":
            return "10"
        elif service_type == "esse":
            return "3-5"
        elif service_type == "kurs":
            return "20-30"
        elif service_type == "referat":
            return "5-10"
        else:
            return "10-15"
            
    num = int(match.group(1))
    
    if service_type == "mustaqil":
        if num <= 15: return "10-15"
        elif num <= 20: return "15-20"
        elif num <= 25: return "20-25"
        else: return "25-30"
    elif service_type == "referat":
        if num <= 10: return "5-10"
        elif num <= 20: return "15-20"
        elif num <= 25: return "20-25"
        else: return "25-30"
    elif service_type == "taqdimot":
        if num <= 12: return "10"
        else: return "15"
    elif service_type == "esse":
        return "3-5"
    elif service_type == "kurs":
        if num <= 30: return "20-30"
        else: return "31-50"
    return "10-15"

async def handle_order_payment(bot, user_id: int, service_type: str, tier: str, amount: float, params_dict: dict, event_source):
    """Processes order payment: checks free usage first, then balance, then falls back to Click."""
    import uuid
    import json
    from database import check_free_usage, mark_free_usage, create_order, get_balance, update_balance, update_order_status
    from services.generation_service import fulfill_order
    from keyboards.inline_keyboards import payment_confirm_kb, main_menu_kb
    from config import CLICK_CARD_NUMBER
    
    is_free = await check_free_usage(user_id, service_type)
    if is_free:
        await mark_free_usage(user_id, service_type)
        mock_order = {
            "order_id": "FREE_USAGE",
            "user_id": user_id,
            "service_type": service_type,
            "pages": tier,
            "parameters": json.dumps(params_dict)
        }
        status_msg = (
            f"🎁 <b>Sizda bepul foydalanish mavjud!</b>\n\n"
            f"Mavzu: {params_dict['topic']}\n"
            f"Hajmi: {tier}\n\n"
            "⏳ Ish tayyorlanmoqda, iltimos kuting..."
        )
        if hasattr(event_source, "edit_text"):
            await event_source.edit_text(status_msg, parse_mode="HTML")
        else:
            await event_source.answer(status_msg, parse_mode="HTML")
            
        await fulfill_order(bot, mock_order)
        if not hasattr(event_source, "edit_text"):
            await event_source.answer("Asosiy menyu:", reply_markup=main_menu_kb())
        return

    order_id = f"ORDER_{uuid.uuid4().hex[:8].upper()}"
    await create_order(order_id, user_id, service_type, tier, amount, json.dumps(params_dict))
    
    balance = await get_balance(user_id)
    if balance >= amount:
        await update_balance(user_id, -amount)
        await update_order_status(order_id, 'paid')
        
        order = {
            "order_id": order_id,
            "user_id": user_id,
            "service_type": service_type,
            "pages": tier,
            "amount": amount,
            "parameters": json.dumps(params_dict),
            "status": "paid"
        }
        
        success_msg = f"💰 To'lov {amount:,} so'm balansingizdan yechildi!\n\n⏳ Hujjat tayyorlanmoqda..."
        if hasattr(event_source, "edit_text"):
            await event_source.edit_text(success_msg)
        else:
            await event_source.answer(success_msg)
            
        await fulfill_order(bot, order)
    else:
        text = (
            f"📄 <b>Mavzu:</b> {params_dict['topic']}\n"
            f"💰 <b>Narx:</b> {amount:,} so'm\n"
            f"⚠️ Balansingiz ({balance:,} so'm) yetarli emas.\n\n"
            f"💳 <b>Karta orqali to'lov (Click)</b>\n"
            f"Quyidagi karta raqamiga to'lovni amalga oshiring:\n"
            f"<code>{CLICK_CARD_NUMBER}</code>\n\n"
            f"To'lovdan so'ng <b>\"✅ To'ladim\"</b> tugmasini bosing va to'lov chekini yuboring."
        )
        if hasattr(event_source, "edit_text"):
            await event_source.edit_text(text, parse_mode="HTML", reply_markup=payment_confirm_kb(order_id))
        else:
            await event_source.answer(text, parse_mode="HTML", reply_markup=payment_confirm_kb(order_id))
